# /// script
# requires-python = ">=3.12"
# dependencies = [
#   "python-pptx==1.0.2",
#   "Pillow==12.3.0",
# ]
# ///
"""Generate the frozen Gloss comparative-v1 repository baselines.

These are deterministic generation-path baselines, not model submissions. Each
path starts from the public prompts and assets and emits a new editable PPTX.
The three runs differ only by their published seed and bounded layout jitter.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import random
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Final

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Cm, Pt

SLIDE_WIDTH_EMU: Final = 12_192_000
SLIDE_HEIGHT_EMU: Final = 6_858_000
SLIDE_WIDTH_CM: Final = 33.8666667
COLORS: Final = {
    "navy": "1B2A4A",
    "coral": "E8634A",
    "teal": "2AACB8",
    "gold": "D4A843",
    "dark": "0F1923",
    "panel": "161F2E",
    "paper": "F5F3EE",
    "white": "FFFFFF",
    "charcoal": "2D2D2D",
}
ALLOWED_FONTS: Final = {
    "title": "Carlito",
    "body": "Liberation Sans",
    "mono": "Liberation Mono",
    "serif": "Liberation Serif",
    "arabic": "Noto Sans Arabic",
    "cjk": "Noto Sans CJK JP",
}


@dataclass(frozen=True)
class PathSpec:
    key: str
    label: str
    prompt_variant: str
    native_objects: bool
    precise_geometry: bool
    full_detail: bool
    jitter_cm: float


PATHS: Final = {
    "native-precise": PathSpec(
        "native-precise",
        "Native precise",
        "canonical",
        native_objects=True,
        precise_geometry=True,
        full_detail=True,
        jitter_cm=0.02,
    ),
    "native-fast": PathSpec(
        "native-fast",
        "Native fast",
        "canonical",
        native_objects=True,
        precise_geometry=False,
        full_detail=False,
        jitter_cm=0.14,
    ),
    "visual-precise": PathSpec(
        "visual-precise",
        "Visual precise",
        "paraphrase-a",
        native_objects=False,
        precise_geometry=True,
        full_detail=True,
        jitter_cm=0.03,
    ),
    "visual-fast": PathSpec(
        "visual-fast",
        "Visual fast",
        "paraphrase-a",
        native_objects=False,
        precise_geometry=False,
        full_detail=False,
        jitter_cm=0.18,
    ),
}
RUN_SEEDS: Final = {1: 1103, 2: 2207, 3: 3301}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_alpha(shape: object, opacity: float) -> None:
    """Set DrawingML alpha after assigning a solid fill."""
    solid_fill = shape.fill._xPr.solidFill  # type: ignore[attr-defined]
    color = next(iter(solid_fill))
    for alpha in color.findall(qn("a:alpha")):
        color.remove(alpha)
    color.append(
        parse_xml(f'<a:alpha {nsdecls("a")} val="{round(opacity * 100000)}"/>')
    )


def set_rtl(paragraph: object) -> None:
    p_pr = paragraph._p.get_or_add_pPr()  # type: ignore[attr-defined]
    p_pr.set("rtl", "1")


def add_slide_number_field(paragraph: object, fallback: str) -> None:
    field = parse_xml(
        f'<a:fld {nsdecls("a")} id="{{{str(uuid.uuid4()).upper()}}}" type="slidenum">'
        '<a:rPr lang="en-US" dirty="0"/><a:t>'
        f"{fallback}</a:t></a:fld>"
    )
    paragraph._p.append(field)  # type: ignore[attr-defined]


def set_gradient(shape: object, left: str, right: str) -> None:
    fill = shape.fill._xPr  # type: ignore[attr-defined]
    insertion_index = len(fill)
    for child in list(fill):
        if child.tag in {qn("a:solidFill"), qn("a:noFill"), qn("a:gradFill")}:
            insertion_index = min(insertion_index, fill.index(child))
            fill.remove(child)
    fill.insert(
        insertion_index,
        parse_xml(
            f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>'
            f'<a:gs pos="0"><a:srgbClr val="{left}"/></a:gs>'
            f'<a:gs pos="100000"><a:srgbClr val="{right}"/></a:gs>'
            '</a:gsLst><a:lin ang="0" scaled="1"/></a:gradFill>'
        ),
    )


def normalize_package(path: Path) -> None:
    """Remove template printer metadata prohibited by the benchmark profile."""
    normalized = path.with_suffix(".normalized.pptx")
    with zipfile.ZipFile(path) as source:
        parts = {
            name: source.read(name)
            for name in source.namelist()
            if not name.startswith("ppt/printerSettings/")
        }

    content_types = etree.fromstring(parts["[Content_Types].xml"])
    for child in list(content_types):
        if (
            child.get("Extension") == "bin"
            and child.get("ContentType")
            == "application/vnd.openxmlformats-officedocument.presentationml.printerSettings"
        ):
            content_types.remove(child)
    parts["[Content_Types].xml"] = etree.tostring(
        content_types,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    relationships_path = "ppt/_rels/presentation.xml.rels"
    relationships = etree.fromstring(parts[relationships_path])
    for child in list(relationships):
        if str(child.get("Type", "")).endswith("/printerSettings"):
            relationships.remove(child)
    parts[relationships_path] = etree.tostring(
        relationships,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
    )

    chart_names = sorted(name for name in parts if name.startswith("ppt/charts/chart"))
    for chart_index, name in enumerate(chart_names, start=1):
        chart = etree.fromstring(parts[name])
        axis_values: dict[str, str] = {}
        for axis in chart.xpath(
            ".//c:axId | .//c:crossAx",
            namespaces={"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"},
        ):
            old_value = str(axis.get("val"))
            if old_value not in axis_values:
                axis_values[old_value] = str(
                    2_000_000_000 + chart_index * 10 + len(axis_values)
                )
            axis.set("val", axis_values[old_value])
        parts[name] = etree.tostring(
            chart,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )

    with zipfile.ZipFile(
        normalized,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
    ) as package:
        for name, data in sorted(parts.items()):
            info = zipfile.ZipInfo(name, date_time=(2025, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            package.writestr(info, data)
    normalized.replace(path)


class DeckBuilder:
    def __init__(self, spec: PathSpec, run: int, seed: int, root: Path) -> None:
        self.spec = spec
        self.run = run
        self.seed = seed
        self.root = root
        self.rng = random.Random(seed)
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH_EMU
        self.prs.slide_height = SLIDE_HEIGHT_EMU
        self.prs.core_properties.title = (
            f"Gloss comparative v1 — {spec.label} — run {run}"
        )
        self.prs.core_properties.subject = (
            f"deterministic generation-path baseline seed {seed}"
        )
        self.prs.core_properties.author = "Gloss contributors"
        self.prs.core_properties.keywords = (
            "Gloss, Gloss, deterministic baseline, editable PowerPoint"
        )
        self.assets = root / "assets" / "mirrored"

    def j(self, value: float) -> float:
        approximation_bias = 0.0 if self.spec.precise_geometry else 0.10
        return (
            value
            + approximation_bias
            + self.rng.uniform(
                -self.spec.jitter_cm,
                self.spec.jitter_cm,
            )
        )

    def add_slide(self, layout: int = 6, title: str | None = None) -> object:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = rgb(COLORS["dark"])
        if title is not None:
            if slide.shapes.title is not None:
                title_shape = slide.shapes.title
                title_shape.left = Cm(self.j(1.5))
                title_shape.top = Cm(self.j(0.55))
                title_shape.width = Cm(30)
                title_shape.height = Cm(1.5)
                title_shape.text = title
                self.style_text(title_shape.text_frame, 36, COLORS["white"], bold=True)
            else:
                self.text(slide, title, 1.5, 0.55, 30, 1.5, 36, bold=True, font="title")
        self.footer(slide, len(self.prs.slides))
        return slide

    def footer(self, slide: object, number: int) -> None:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Cm(0),
            Cm(17.5),
            SLIDE_WIDTH_EMU,
            Cm(17.5),
        )
        line.line.color.rgb = rgb(COLORS["teal"])
        line.line.width = Pt(0.5)
        self.text(slide, "GLOSS", 0.6, 17.62, 4, 0.45, 8, color=COLORS["white"])
        number_box = self.text(
            slide,
            "",
            31.8,
            17.55,
            1.2,
            0.55,
            10,
            color=COLORS["white"],
            align=PP_ALIGN.RIGHT,
        )
        add_slide_number_field(number_box.text_frame.paragraphs[0], str(number))

    def style_text(
        self,
        text_frame: object,
        size: float,
        color: str,
        *,
        bold: bool = False,
        font: str = "body",
        align: PP_ALIGN | None = None,
        italic: bool = False,
    ) -> None:
        text_frame.word_wrap = True
        text_frame.margin_left = Cm(0.12)
        text_frame.margin_right = Cm(0.12)
        text_frame.margin_top = Cm(0.06)
        text_frame.margin_bottom = Cm(0.06)
        for paragraph in text_frame.paragraphs:
            if align is not None:
                paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = ALLOWED_FONTS[font]
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = rgb(color)

    def text(
        self,
        slide: object,
        value: str,
        x: float,
        y: float,
        w: float,
        h: float,
        size: float = 18,
        *,
        color: str = COLORS["white"],
        bold: bool = False,
        font: str = "body",
        align: PP_ALIGN = PP_ALIGN.LEFT,
        valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
        rotation: float = 0,
        italic: bool = False,
    ) -> object:
        shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
        shape.rotation = rotation
        shape.text_frame.text = value
        shape.text_frame.vertical_anchor = valign
        self.style_text(
            shape.text_frame,
            size,
            color,
            bold=bold,
            font=font,
            align=align,
            italic=italic,
        )
        return shape

    def box(
        self,
        slide: object,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: str,
        *,
        text: str = "",
        opacity: float = 1,
        rounded: bool = True,
        line: str | None = None,
        size: float = 14,
        text_color: str = COLORS["white"],
        rotation: float = 0,
    ) -> object:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Cm(x), Cm(y), Cm(w), Cm(h))
        shape.rotation = rotation
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if opacity < 1:
            add_alpha(shape, opacity)
        if line:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        if text:
            shape.text = text
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            self.style_text(
                shape.text_frame,
                size,
                text_color,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        return shape

    def picture(
        self,
        slide: object,
        asset: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        crop_left: float = 0,
        crop_top: float = 0,
        crop_right: float = 0,
        crop_bottom: float = 0,
    ) -> object:
        shape = slide.shapes.add_picture(
            str(self.assets / asset), Cm(x), Cm(y), Cm(w), Cm(h)
        )
        shape.crop_left = crop_left
        shape.crop_top = crop_top
        shape.crop_right = crop_right
        shape.crop_bottom = crop_bottom
        return shape

    def connector(
        self,
        slide: object,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        color: str = COLORS["paper"],
        width: float = 1.5,
    ) -> object:
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2)
        )
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)
        return shape

    def table(
        self,
        slide: object,
        rows: list[list[str]],
        x: float,
        y: float,
        w: float,
        h: float,
    ) -> object:
        table = slide.shapes.add_table(
            len(rows), len(rows[0]), Cm(x), Cm(y), Cm(w), Cm(h)
        ).table
        for row_idx, values in enumerate(rows):
            for col_idx, value in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(
                    COLORS["navy"]
                    if row_idx == 0
                    else COLORS["dark" if row_idx % 2 else "panel"]
                )
                self.style_text(
                    cell.text_frame,
                    12 if row_idx else 13,
                    COLORS["white" if row_idx == 0 else "paper"],
                    bold=row_idx == 0,
                    align=PP_ALIGN.CENTER if col_idx else PP_ALIGN.LEFT,
                )
        return table

    def chart(
        self,
        slide: object,
        chart_type: XL_CHART_TYPE,
        categories: list[str],
        series: list[tuple[str, list[float]]],
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        title: str = "",
    ) -> object:
        data = CategoryChartData()
        data.categories = categories
        for name, values in series:
            data.add_series(name, values)
        chart = slide.shapes.add_chart(
            chart_type, Cm(x), Cm(y), Cm(w), Cm(h), data
        ).chart
        chart.has_title = bool(title)
        if title:
            chart.chart_title.text_frame.text = title
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        for index, chart_series in enumerate(chart.series):
            chart_series.format.fill.solid()
            chart_series.format.fill.fore_color.rgb = rgb(
                [COLORS["navy"], COLORS["coral"], COLORS["teal"]][index % 3]
            )
        return chart

    def manual_bar_chart(
        self,
        slide: object,
        categories: list[str],
        values: list[float],
        x: float,
        y: float,
        w: float,
        h: float,
    ) -> None:
        maximum = max(values)
        row_h = h / len(values)
        for index, (category, value) in enumerate(zip(categories, values, strict=True)):
            row_y = y + index * row_h
            self.text(slide, category, x, row_y, w * 0.34, row_h, 10)
            self.box(
                slide,
                x + w * 0.36,
                row_y + row_h * 0.2,
                w * 0.55 * value / maximum,
                row_h * 0.56,
                COLORS["coral"],
                rounded=False,
            )
            self.text(slide, f"{value:g}", x + w * 0.92, row_y, w * 0.08, row_h, 10)

    def build(self) -> None:
        self.slide_01()
        self.slide_02()
        self.slide_03()
        self.slide_04()
        self.slide_05()
        self.slide_06()
        self.slide_07()
        self.slide_08()
        self.slide_09()
        self.slide_10()
        self.slide_11()
        self.slide_12()
        self.slide_13()
        self.slide_14()
        self.slide_15()
        self.slide_16()
        self.slide_17()
        self.slide_18()
        self.slide_19()
        self.slide_20()

    def slide_01(self) -> None:
        slide = self.add_slide(0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 3:
                placeholder.text = "Gloss v1"
                self.style_text(
                    placeholder.text_frame,
                    44,
                    COLORS["white"],
                    bold=True,
                    font="title",
                    align=PP_ALIGN.CENTER,
                )
            elif placeholder.placeholder_format.type == 4:
                placeholder.text = "Benchmark for Slide Generation Fidelity"
                self.style_text(
                    placeholder.text_frame,
                    24,
                    COLORS["coral"],
                    font="title",
                    align=PP_ALIGN.CENTER,
                )
        self.box(slide, 16, 8, 10, 6, COLORS["gold"], opacity=0.25)
        self.box(slide, 14, 4, 8, 8, COLORS["teal"], opacity=0.40)
        self.picture(slide, "hero-abstract.png", 18.867, 2, 15, 10.2, crop_left=0.2)
        self.box(slide, 18, 1, 6, 10, COLORS["coral"], opacity=0.30)

    def slide_02(self) -> None:
        slide = self.add_slide(1, "Agenda")
        body = next((p for p in slide.placeholders if p != slide.shapes.title), None)
        agenda = [
            ("Opening Remarks", 0),
            ("Welcome and introductions", 1),
            ("Safety briefing", 1),
            ("Technical Deep Dive", 0),
            ("Architecture overview", 1),
            ("Performance benchmarks", 1),
            ("Security audit results", 1),
            ("Breakout Sessions", 0),
            ("Track A: Infrastructure", 1),
            ("Track B: Applications", 1),
            ("Closing & Next Steps", 0),
        ]
        if body is not None:
            body.left, body.top, body.width, body.height = (
                Cm(1.5),
                Cm(3),
                Cm(20),
                Cm(13.5),
            )
            body.text_frame.clear()
            for index, (value, level) in enumerate(agenda):
                paragraph = (
                    body.text_frame.paragraphs[0]
                    if index == 0
                    else body.text_frame.add_paragraph()
                )
                paragraph.text = value
                paragraph.level = level
                for run in paragraph.runs:
                    run.font.name = ALLOWED_FONTS["body"]
                    run.font.size = Pt(18 if level == 0 else 14)
                    run.font.bold = level == 0
                    run.font.color.rgb = rgb(
                        COLORS["coral"] if level == 0 else COLORS["paper"]
                    )
        for index, label in enumerate(("09:00", "11:30", "14:00")):
            y = 4.2 + index * 4.1
            self.box(slide, 24, y, 1, 1, COLORS["teal"], rounded=True)
            self.text(slide, label, 25.4, y - 0.1, 3, 1.2, 14, bold=True)

    def slide_03(self) -> None:
        slide = self.add_slide(6, "Performance Metrics")
        rows = [
            ["Metric", "Q1 2024", "Q2 2024", "Q3 2024", "Target"],
            ["Latency (p50)", "12ms", "11ms", "9ms", "≤10ms"],
            ["Latency (p99)", "89ms", "72ms", "58ms", "≤75ms"],
            ["Throughput", "1,200 rps", "1,450 rps", "1,890 rps", "≥2,000 rps"],
            ["Error Rate", "0.12%", "0.08%", "0.04%", "≤0.05%"],
            ["Uptime", "99.91%", "99.95%", "99.98%", "≥99.95%"],
            ["Cache Hit", "72%", "78%", "84%", "≥80%"],
        ]
        if not self.spec.full_detail:
            rows = rows[:-2]
        if self.spec.native_objects:
            self.table(slide, rows, 1.5, 3, 30, 12)
        else:
            for row_index, row in enumerate(rows):
                for col_index, value in enumerate(row):
                    self.box(
                        slide,
                        1.5 + col_index * 6,
                        3 + row_index * 1.7,
                        6,
                        1.7,
                        COLORS["navy"] if row_index == 0 else COLORS["panel"],
                        text=value,
                        rounded=False,
                        line=COLORS["teal"],
                        size=11,
                    )
        self.box(
            slide,
            24,
            15.2,
            8,
            2.2,
            COLORS["coral"],
            text="Throughput target\nmissed by 5.5%",
            size=12,
        )

    def slide_04(self) -> None:
        slide = self.add_slide(1, "Revenue Growth")
        categories = [
            "North America",
            "Europe",
            "Asia Pacific",
            "Latin America",
            "Middle East & Africa",
        ]
        if self.spec.native_objects:
            self.chart(
                slide,
                XL_CHART_TYPE.BAR_CLUSTERED,
                categories,
                [
                    ("2023", [42.3, 28.7, 19.4, 8.1, 3.9]),
                    ("2024", [51.8, 35.2, 31.6, 12.4, 7.2]),
                ],
                1.5,
                3,
                20,
                13.5,
                title="Annual Revenue by Region",
            )
        else:
            self.manual_bar_chart(
                slide, categories, [51.8, 35.2, 31.6, 12.4, 7.2], 1.5, 3.5, 20, 11
            )
        self.text(
            slide,
            "Key Takeaways:\n\n▸ Total revenue up 38.2% YoY\n▸ APAC fastest-growing region\n▸ All regions exceeded targets",
            23,
            4,
            9,
            7,
            14,
        )
        self.box(slide, 16, 8.2, 5, 2, COLORS["teal"], text="+62.9% YoY")
        self.box(slide, 16, 12.6, 5, 2, COLORS["gold"], text="+84.6% YoY")

    def slide_05(self) -> None:
        slide = self.add_slide(1, "Our Team")
        people = [
            (
                "Sarah Chen",
                "Chief Architect",
                "15 years building distributed systems at scale",
            ),
            (
                "Marcus Rivera",
                "Head of Design",
                "Formerly at Apple, obsessed with pixel-perfect interfaces",
            ),
            (
                "Yuki Tanaka",
                "ML Engineering Lead",
                "PhD in NLP, built the model serving infrastructure",
            ),
        ]
        for index, (name, role, description) in enumerate(people):
            x = 1.5 + index * 10.5
            self.box(slide, x, 3.1, 9, 12.8, COLORS["panel"], line=COLORS["teal"])
            self.box(
                slide, x + 2.5, 4, 4, 4, COLORS["navy"], text="Photo", rounded=True
            )
            self.text(
                slide,
                name,
                x + 0.6,
                8.5,
                7.8,
                1,
                18,
                bold=True,
                font="title",
                align=PP_ALIGN.CENTER,
            )
            self.text(
                slide,
                role,
                x + 0.6,
                9.6,
                7.8,
                1,
                14,
                color=COLORS["coral"],
                align=PP_ALIGN.CENTER,
            )
            self.text(
                slide, description, x + 0.7, 11, 7.6, 2.8, 12, align=PP_ALIGN.CENTER
            )

    def slide_06(self) -> None:
        slide = self.add_slide(6)
        columns = [
            (
                "Global Perspectives",
                "The rapid advancement of language models\nhas transformed how organizations approach\ndocument generation.",
                "body",
                False,
            ),
            (
                "وجهات نظر عالمية",
                "لقد أدى التقدم السريع في نماذج اللغة إلى تحويل\nكيفية تعامل المؤسسات مع إنشاء المستندات.",
                "arabic",
                True,
            ),
            (
                "グローバルな視点",
                "言語モデルの急速な進歩により、組織がドキュメント生成に\n取り組む方法が変革されました。",
                "cjk",
                False,
            ),
        ]
        for index, (heading, body, font, rtl) in enumerate(columns):
            x = 1 + index * 11
            self.text(
                slide,
                heading,
                x,
                2,
                10,
                1.2,
                24,
                bold=True,
                font=font,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT,
            )
            body_shape = self.text(
                slide,
                body,
                x,
                3.5,
                10,
                9,
                14,
                font=font,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.TOP,
            )
            if rtl:
                for paragraph in body_shape.text_frame.paragraphs:
                    set_rtl(paragraph)
        self.connector(slide, 11.2, 3, 11.2, 16, color=COLORS["teal"], width=1)
        self.connector(slide, 22.5, 3, 22.5, 16, color=COLORS["teal"], width=1)
        self.box(
            slide,
            15,
            13,
            12,
            4,
            COLORS["coral"],
            text="AI-generated slides must handle RTL text,\nCJK line breaking, and mixed scripts correctly.",
            opacity=0.9,
            size=13,
        )

    def slide_07(self) -> None:
        slide = self.add_slide(6, "Image Handling")
        self.picture(
            slide, "cityscape.png", 1.5, 2.5, 10, 8, crop_top=0.15, crop_bottom=0.10
        )
        self.picture(slide, "cityscape.png", 13, 2.5, 8, 8, crop_left=0.30)
        self.picture(slide, "texture-pattern.png", 22.5, 2.5, 10, 8, crop_right=0.25)
        for x, label in zip(
            (1.5, 12, 22.5),
            ("Full panorama crop", "Circle-masked detail", "Rounded corner crop"),
            strict=True,
        ):
            self.text(
                slide, label, x, 10.8, 10, 0.8, 12, italic=True, align=PP_ALIGN.CENTER
            )
        self.box(
            slide, 14, 3.5, 6, 6, COLORS["navy"], text="PREVIEW", opacity=0.5, size=24
        )

    def slide_08(self) -> None:
        slide = self.add_slide(6, "Depth & Layering")
        cards = [
            (4, 6, COLORS["navy"], 1.0, COLORS["white"]),
            (7, 5, COLORS["teal"], 0.85, COLORS["white"]),
            (10, 4, COLORS["coral"], 0.70, COLORS["white"]),
            (13, 3, COLORS["gold"], 0.55, COLORS["navy"]),
            (16, 2, COLORS["white"], 0.90, COLORS["navy"]),
        ]
        for index, (x, y, fill, opacity, text_color) in enumerate(cards, start=1):
            self.box(
                slide,
                x,
                y,
                12,
                8,
                fill,
                text=f"Layer {index}",
                opacity=opacity,
                size=20,
                text_color=text_color,
            )
        gradient = self.box(
            slide, 0, 15, SLIDE_WIDTH_CM, 2.5, COLORS["dark"], rounded=False
        )
        set_gradient(gradient, COLORS["dark"], COLORS["teal"])

    def slide_09(self) -> None:
        slide = self.add_slide(1, "API Reference")
        core = (
            "GET /api/v1/submissions\nAuthorization: Bearer YOUR_TOKEN_HERE\nContent-Type: application/json\n\n"
            "Query Parameters:\n  model_id    string  required  Filter by model\n  tier        integer optional  1, 2, or 3\n"
            "  status      string  optional  queued|grading|completed\n  limit       integer optional  Default: 20, max: 100\n\n"
            'Response 200:\n{\n  "submissions": [{"id": "sub_abc123", "model_id": "gpt-4o", "tier": 3, '
            '"status": "completed", "fidelity_score": 0.847}],\n  "total": 142,\n  "has_more": true\n}'
        )
        extra = "\n\nError Responses:\n  401 Unauthorized\n  403 Forbidden\n  404 Not Found\n\nRate Limits:\n  Free tier: 10 req/min\n  Pro tier: 60 req/min"
        self.text(
            slide,
            "Fixed Size (clips)",
            1.5,
            2.5,
            14,
            0.6,
            11,
            color=COLORS["coral"],
            bold=True,
        )
        left = self.text(
            slide, core, 1.5, 3.2, 14, 13.7, 11, font="mono", valign=MSO_ANCHOR.TOP
        )
        left.text_frame.auto_size = None
        self.text(
            slide,
            "Auto-Shrink (fits)",
            17.2,
            2.5,
            14,
            0.6,
            11,
            color=COLORS["teal"],
            bold=True,
        )
        right = self.text(
            slide,
            core + extra,
            17.2,
            3.2,
            14,
            13.7,
            11,
            font="mono",
            valign=MSO_ANCHOR.TOP,
        )
        right.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def slide_10(self) -> None:
        slide = self.add_slide(6, "System Architecture")
        boxes = {
            "Client App": (2, 3, COLORS["teal"]),
            "API Gateway": (14, 3, COLORS["navy"]),
            "Auth Service": (8, 9, COLORS["coral"]),
            "Job Queue": (14, 9, COLORS["gold"]),
            "Grader": (20, 9, COLORS["coral"]),
            "PostgreSQL": (6, 14.2, COLORS["navy"]),
            "Object Store": (14, 14.2, COLORS["navy"]),
            "Leaderboard": (22, 14.2, COLORS["teal"]),
        }
        for label, (x, y, fill) in boxes.items():
            self.box(slide, x, y, 5, 2.4, fill, text=label, size=12)
        links = [
            (7, 4.2, 14, 4.2),
            (16.5, 5.4, 10.5, 9),
            (16.5, 5.4, 16.5, 9),
            (19, 10.2, 20, 10.2),
            (22.5, 11.4, 16.5, 14.2),
            (14, 15.4, 11, 15.4),
            (19, 15.4, 22, 15.4),
        ]
        for coords in links:
            self.connector(slide, *coords)
        self.text(slide, "HTTPS", 9.2, 3.4, 2, 0.6, 10)
        self.text(slide, "JWT verify", 11.5, 6.5, 2.5, 0.6, 10)
        self.text(slide, "Enqueue", 16.8, 6.8, 2.2, 0.6, 10)

    def slide_11(self) -> None:
        slide = self.add_slide(1, "Brand Colors")
        swatches = [
            ("Primary", COLORS["navy"]),
            ("Secondary", COLORS["coral"]),
            ("Accent 1", COLORS["teal"]),
            ("Override", "FF6B35"),
            ("Override", "7B2D8E"),
            ("Override", "3D9970"),
        ]
        for index, (label, fill) in enumerate(swatches):
            row, col = divmod(index, 3)
            x, y = 5 + col * 9, 3.2 + row * 6
            self.box(slide, x, y, 4, 4, fill)
            self.text(
                slide,
                label,
                x - 0.5,
                y + 4.1,
                5,
                0.7,
                12,
                color=fill if row else COLORS["paper"],
                align=PP_ALIGN.CENTER,
            )
        self.text(
            slide,
            "Top row uses theme roles. Bottom row uses explicit RGB values.\nChanging the theme updates the top row but not the bottom row.",
            4.5,
            15.6,
            25,
            1.2,
            13,
            align=PP_ALIGN.CENTER,
        )

    def slide_12(self) -> None:
        slide = self.add_slide(1, "Document Fields")
        field = self.text(slide, "Current slide: ", 3, 4, 12, 1.4, 18)
        add_slide_number_field(field.text_frame.paragraphs[0], "12")
        self.text(
            slide,
            "← Live slide number field (not static text)",
            15.5,
            4,
            13,
            1.4,
            12,
            color=COLORS["teal"],
            italic=True,
        )
        self.text(slide, "Generated: 2025-01-15 10:30 UTC", 3, 7, 12, 1.4, 18)
        self.text(
            slide,
            "← Fixed date/time field",
            15.5,
            7,
            13,
            1.4,
            12,
            color=COLORS["teal"],
            italic=True,
        )
        self.text(slide, "Gloss Benchmark v1", 3, 10, 12, 1.4, 18)
        self.text(
            slide,
            "← Footer from master placeholder",
            15.5,
            10,
            13,
            1.4,
            12,
            color=COLORS["teal"],
            italic=True,
        )
        self.text(slide, "Slide 12", 3, 13.5, 12, 1.4, 18)

    def slide_13(self) -> None:
        slide = self.add_slide(6, "Composite Stress")
        categories = ["Structure", "Text", "Visual", "Assets", "Fields", "Other"]
        if self.spec.native_objects:
            data = ChartData()
            data.categories = categories
            data.add_series("Checks", (28, 22, 18, 14, 10, 8))
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, Cm(1.5), Cm(3), Cm(10), Cm(8), data
            ).chart
            chart.has_legend = False
            chart.plots[0].has_data_labels = True
            chart.plots[0].data_labels.show_percentage = True
        else:
            self.manual_bar_chart(
                slide, categories, [28, 22, 18, 14, 10, 8], 1.5, 3, 10, 8
            )
        rows = [
            ["Check", "Status", "Weight"],
            ["Schema", "Pass", "3"],
            ["Rendering", "Pass", "2"],
            ["Assets", "Review", "1"],
        ]
        if self.spec.native_objects:
            self.table(slide, rows, 13, 3, 9, 7)
        else:
            for row_index, row in enumerate(rows):
                for col_index, value in enumerate(row):
                    self.box(
                        slide,
                        13 + col_index * 3,
                        3 + row_index * 1.7,
                        3,
                        1.7,
                        COLORS["panel"],
                        text=value,
                        rounded=False,
                        size=10,
                    )
        self.picture(slide, "cityscape.png", 23, 3, 8, 6, crop_left=0.15)
        arabic = self.text(
            slide,
            "الهيكل والمحتوى والعرض في اختبار واحد",
            12,
            12.5,
            16,
            1.5,
            14,
            font="arabic",
            align=PP_ALIGN.RIGHT,
        )
        set_rtl(arabic.text_frame.paragraphs[0])
        for x, label in (
            (2, "NATIVE CHART"),
            (13, "NATIVE TABLE"),
            (23, "APPROVED ASSET"),
        ):
            self.box(slide, x, 10, 7, 1.4, COLORS["teal"], text=label, size=11)

    def slide_14(self) -> None:
        slide = self.add_slide(3, "RTL Systems Review")
        arabic_lines = [
            "مراجعة الأنظمة الموزعة",
            "دقة العرض",
            "سلامة البنية",
            "The الذكاء الاصطناعي revolution",
            "الإصدار Gloss v1",
        ]
        left = self.text(
            slide,
            "\n\n".join(arabic_lines),
            1.5,
            3,
            14.6,
            13,
            14,
            font="arabic",
            align=PP_ALIGN.RIGHT,
            valign=MSO_ANCHOR.TOP,
        )
        for paragraph in left.text_frame.paragraphs:
            set_rtl(paragraph)
        english = [
            "A structural benchmark must preserve meaning and direction.",
            "Rendering alone cannot establish editability.",
            "Native objects keep content reusable.",
            "Bidirectional text must retain Unicode order.",
            "Every claim needs a reproducible artifact.",
        ]
        self.text(
            slide, "\n\n".join(english), 17.7, 3, 14.6, 13, 14, valign=MSO_ANCHOR.TOP
        )
        self.connector(slide, 16.93, 3, 16.93, 16, color=COLORS["teal"], width=1)

    def slide_15(self) -> None:
        slide = self.add_slide(6, "Rotation Atlas")
        entries = [
            ("Baseline 0°", 0),
            ("Diagonal 45°", 45),
            ("Vertical 90°", 90),
            ("Reverse 135°", 135),
            ("Vertical 270°", 270),
        ]
        for index, (label, angle) in enumerate(entries):
            x = 1.5 + index * 6.2
            self.box(slide, x, 7.5, 6, 2, COLORS["navy"], rotation=angle)
            self.text(
                slide,
                label,
                x + 0.5,
                7.75,
                5,
                1.5,
                16,
                bold=True,
                align=PP_ALIGN.CENTER,
                rotation=angle,
            )

    def slide_16(self) -> None:
        slide = self.add_slide(6, "Beyond the Frame")
        self.box(slide, -3, -2, 10, 10, COLORS["coral"], rounded=True)
        self.box(slide, 28.867, 6, 12, 6, COLORS["teal"], rounded=False)
        self.box(slide, 8, -1, 16, 3, COLORS["gold"], rounded=False)
        self.text(
            slide,
            "INTENTIONAL BLEED",
            8,
            7,
            18,
            2,
            28,
            bold=True,
            font="title",
            align=PP_ALIGN.CENTER,
        )
        self.text(
            slide,
            "Negative coordinates are part of the composition.",
            7,
            10,
            20,
            1.5,
            18,
            align=PP_ALIGN.CENTER,
        )

    def slide_17(self) -> None:
        slide = self.add_slide(6, "Nested Systems")
        labels = ("A1", "A2", "B1", "B2", "C1", "C2")
        for index, label in enumerate(labels):
            col, row = divmod(index, 2)
            x, y = 3.8 + col * 9, 4.2 + row * 6
            self.box(slide, x, y, 2.2, 2.2, COLORS["teal"], rounded=True)
            self.box(
                slide, x + 1.1, y + 0.8, 5.3, 2.6, COLORS["navy"], line=COLORS["coral"]
            )
            self.text(
                slide, label, x + 2, y + 1.4, 3, 1, 14, bold=True, align=PP_ALIGN.CENTER
            )
        for index, label in enumerate(("INPUT", "PROCESS", "OUTPUT")):
            self.text(
                slide,
                label,
                4 + index * 9,
                3,
                7,
                1,
                14,
                color=COLORS["gold"],
                bold=True,
                align=PP_ALIGN.CENTER,
            )

    def slide_18(self) -> None:
        slide = self.add_slide(6, "Three Cities / 三つの都市")
        content = [
            (
                "cityscape.png",
                "Systems",
                "Infrastructure becomes legible when structure and meaning stay together.",
                "Build for reuse.",
            ),
            (
                "texture-pattern.png",
                "システム",
                "生成されたスライドは構造と意味を保持します。",
                "意味を保つ。",
            ),
            (
                "cityscape.png",
                "Shared Futures",
                "Editable artifacts let teams continue the work after generation.",
                "The artifact is the product.",
            ),
        ]
        for index, (asset, heading, body, quote) in enumerate(content):
            x = 1.5 + index * 10.4
            if index == 1:
                self.picture(slide, "texture-pattern.png", x, 2.7, 9.5, 13.6)
            self.picture(slide, asset, x, 3, 9.5, 3.5)
            self.text(
                slide,
                heading,
                x,
                6.8,
                9.5,
                1.2,
                20,
                bold=True,
                font="cjk" if index == 1 else "title",
            )
            self.text(
                slide,
                body,
                x,
                8.2,
                9.5,
                4.5,
                13,
                font="cjk" if index == 1 else "body",
                valign=MSO_ANCHOR.TOP,
            )
            self.text(
                slide,
                f"“{quote}”",
                x,
                13.2,
                9.5,
                2.2,
                14,
                color=COLORS["coral"],
                italic=True,
                align=PP_ALIGN.CENTER,
            )

    def slide_19(self) -> None:
        slide = self.add_slide(1, "Design System Audit")
        labels = ["Typography", "Palette", "Spacing", "Line Weight", "Navigation"]
        for index, label in enumerate(labels):
            y = 3.2 + index * 2.5
            self.text(
                slide,
                label,
                2,
                y,
                8,
                1.2,
                18,
                bold=index == 0,
                font="title" if index == 0 else "body",
            )
            self.connector(
                slide, 11, y + 0.6, 29, y + 0.6, color=COLORS["teal"], width=0.5
            )
        self.text(
            slide,
            "Return to Cover",
            12,
            14.9,
            7,
            1,
            14,
            color=COLORS["coral"],
            bold=True,
        )
        self.text(
            slide, "Meet the Team", 21, 14.9, 7, 1, 14, color=COLORS["teal"], bold=True
        )

    def slide_20(self) -> None:
        slide = self.add_slide(6, "Gloss Synthesis")
        categories = ["L1", "L2", "L3"]
        if self.spec.native_objects:
            self.chart(
                slide,
                XL_CHART_TYPE.LINE_MARKERS,
                categories,
                [
                    ("Structural", [72, 84, 96]),
                    ("Visual", [68, 82, 95]),
                    ("Combined", [70, 83, 97]),
                ],
                1,
                2.7,
                10.5,
                5,
            )
        else:
            self.manual_bar_chart(slide, categories, [70, 83, 97], 1, 2.7, 10.5, 5)
        rows = [
            ["Tier", "Slides", "Checks", "SSIM", "Schema", "Status"],
            ["L1", "5", "70", "0.9999", "Pass", "Ready"],
            ["L3", "20", "280", "0.9999", "Pass", "Ready"],
        ]
        if self.spec.native_objects:
            self.table(slide, rows, 12, 2.7, 12, 4)
        else:
            for row_index, row in enumerate(rows):
                for col_index, value in enumerate(row):
                    self.box(
                        slide,
                        12 + col_index * 2,
                        2.7 + row_index * 1.3,
                        2,
                        1.3,
                        COLORS["panel"],
                        text=value,
                        rounded=False,
                        size=8,
                    )
        self.picture(slide, "hero-abstract.png", 25, 2.7, 7.5, 5, crop_left=0.2)
        arabic = self.text(
            slide,
            "اختبار شامل",
            1,
            8.2,
            7,
            1.2,
            18,
            font="arabic",
            align=PP_ALIGN.RIGHT,
        )
        set_rtl(arabic.text_frame.paragraphs[0])
        self.text(slide, "総合テスト", 8.5, 8.2, 7, 1.2, 18, font="cjk")
        self.text(
            slide,
            "ROTATED",
            17,
            8.5,
            5,
            1.5,
            16,
            bold=True,
            rotation=45,
            align=PP_ALIGN.CENTER,
        )
        for index, (fill, opacity) in enumerate(
            ((COLORS["navy"], 0.7), (COLORS["coral"], 0.6), (COLORS["teal"], 0.5))
        ):
            self.box(
                slide, 23 + index * 1.6, 8 + index * 0.8, 6, 4, fill, opacity=opacity
            )
        gradient = self.box(slide, 0, 12.5, 14, 3.8, COLORS["navy"], rounded=False)
        set_gradient(gradient, COLORS["navy"], COLORS["teal"])
        self.text(
            slide,
            "• Looks right\n• Built right\n• Reproduces from source",
            1,
            13,
            12,
            2.6,
            16,
            bold=True,
        )
        field = self.text(slide, "Current slide: ", 16, 14, 9, 1.2, 16)
        add_slide_number_field(field.text_frame.paragraphs[0], "20")


def validate_deck(path: Path) -> None:
    with zipfile.ZipFile(path) as package:
        names = package.namelist()
        slides = [
            name
            for name in names
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        if len(slides) != 20:
            raise ValueError(f"{path}: expected 20 slides, found {len(slides)}")
        if any(name.endswith("vbaProject.bin") for name in names):
            raise ValueError(f"{path}: macros are prohibited")
    presentation = Presentation(path)
    ratio = presentation.slide_width / presentation.slide_height
    if abs(ratio - (16 / 9)) > 0.001:
        raise ValueError(f"{path}: slide ratio is {ratio}")


def build_run(root: Path, output_root: Path, spec: PathSpec, run: int) -> Path:
    seed = RUN_SEEDS[run]
    run_dir = output_root / "runs" / spec.key / f"run-{run}"
    run_dir.mkdir(parents=True, exist_ok=True)
    deck_path = run_dir / "deck.pptx"
    builder = DeckBuilder(spec, run, seed, root)
    builder.build()
    builder.prs.save(deck_path)
    normalize_package(deck_path)
    validate_deck(deck_path)
    generation = {
        "schema_version": "1.0",
        "generator_kind": "deterministic_repository_baseline",
        "verification_label": "local artifact score; self-reported",
        "generation_attribution": "repository-owned path; no model attribution",
        "path_id": spec.key,
        "path_label": spec.label,
        "prompt_variant": spec.prompt_variant,
        "run": run,
        "seed": seed,
        "native_objects": spec.native_objects,
        "precise_geometry": spec.precise_geometry,
        "full_detail": spec.full_detail,
        "source_entrypoint": "generate_baselines.py",
        "generator_sha256": (
            "sha256:" + hashlib.sha256(Path(__file__).read_bytes()).hexdigest()
        ),
        "human_intervention": False,
        "post_processing": False,
        "external_resources_used": False,
    }
    (run_dir / "generation.json").write_text(
        json.dumps(generation, indent=2, sort_keys=True) + "\n", encoding="utf-8"
    )
    return deck_path


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--benchmark-root",
        type=Path,
        default=Path(__file__).resolve().parents[1],
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path(__file__).resolve().parent,
    )
    parser.add_argument("--path", choices=sorted(PATHS))
    parser.add_argument("--run", type=int, choices=(1, 2, 3))
    args = parser.parse_args()
    selected_paths = [PATHS[args.path]] if args.path else list(PATHS.values())
    selected_runs = [args.run] if args.run else [1, 2, 3]
    for spec in selected_paths:
        for run in selected_runs:
            deck = build_run(args.benchmark_root, args.output_root, spec, run)
            print(deck)


if __name__ == "__main__":
    main()
